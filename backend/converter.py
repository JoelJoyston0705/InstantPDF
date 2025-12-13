import os
import shutil
from pathlib import Path
from docx import Document
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
import openpyxl
from fpdf import FPDF, XPos, YPos
import img2pdf
from PIL import Image
from pptx import Presentation
import pdfplumber
from pdf2image import convert_from_path
from pypdf import PdfReader

import subprocess
import platform

def get_libreoffice_command():
    """Find the LibreOffice executable."""
    if platform.system() == "Darwin":  # macOS
        search_paths = [
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            "/usr/local/bin/soffice"
        ]
        for path in search_paths:
            if os.path.exists(path):
                return path
    
    # Default for Linux/Docker (Railway)
    return "soffice"

def convert_with_libreoffice(input_path: str, output_path: str):
    """
    Convert document to PDF using LibreOffice (High Fidelity).
    Returns True if successful, False otherwise.
    """
    soffice = get_libreoffice_command()
    
    # Check if we can run soffice (basic check)
    try:
        if platform.system() != "Darwin":
            subprocess.run([soffice, "--version"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, check=True)
        elif not os.path.exists(soffice):
            return False
    except (subprocess.SubprocessError, FileNotFoundError):
        print("LibreOffice not found, falling back to basic conversion.")
        return False

    out_dir = os.path.dirname(output_path)
    
    # libreoffice --convert-to pdf puts the file in out_dir with the same basename
    # We need to handle the renaming if output_path has a different name
    cmd = [
        soffice,
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        out_dir,
        input_path
    ]
    
    try:
        subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)
        
        # Verify output exists
        # LibreOffice uses the input filename + .pdf
        base_name = os.path.splitext(os.path.basename(input_path))[0]
        expected_output = os.path.join(out_dir, base_name + ".pdf")
        
        if os.path.exists(expected_output):
            # If the requested output path is different, rename it
            if expected_output != output_path:
                if os.path.exists(output_path):
                    os.remove(output_path)
                os.rename(expected_output, output_path)
            return True
    except Exception as e:
        print(f"LibreOffice conversion error: {e}")
        return False
    
    return False

def convert_docx_to_pdf(input_path: str, output_path: str):
    # Try High Fidelity Conversion first
    if convert_with_libreoffice(input_path, output_path):
        return

    try:
        # Load the DOCX document
        doc = Document(input_path)
        
        # Create PDF
        pdf = SimpleDocTemplate(output_path, pagesize=letter)
        story = []
        styles = getSampleStyleSheet()
        
        # Add a custom style for better formatting
        normal_style = styles['Normal']
        
        # Process paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                # Handle different paragraph styles
                if para.style.name.startswith('Heading'):
                    style = styles['Heading1'] if 'Heading 1' in para.style.name else styles['Heading2']
                else:
                    style = normal_style
                
                p = Paragraph(para.text, style)
                story.append(p)
                story.append(Spacer(1, 0.1 * inch))
        
        # If document is empty, add a placeholder
        if not story:
            story.append(Paragraph("Empty document", normal_style))
        
        # Build PDF
        pdf.build(story)
    except Exception as e:
        raise RuntimeError(f"DOCX conversion failed: {e}")

def convert_xlsx_to_pdf(input_path: str, output_path: str):
    try:
        wb = openpyxl.load_workbook(input_path)
        sheet = wb.active
        
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        
        # Add Unicode font support (bundled Roboto)
        current_dir = os.path.dirname(os.path.abspath(__file__))
        font_path = os.path.join(current_dir, "fonts", "Roboto-Regular.ttf")
        font_name = "Roboto"
        use_unicode_font = False
        
        if os.path.exists(font_path):
            try:
                pdf.add_font(font_name, fname=font_path)
                use_unicode_font = True
            except Exception as e:
                print(f"Failed to load font {font_path}: {e}")

        pdf.add_page()
        
        if use_unicode_font:
            pdf.set_font(font_name, size=10)
        else:
            pdf.set_font("Helvetica", size=10)

        # Iterate rows
        for row in sheet.iter_rows(values_only=True):
            for item in row:
                # Basic cell handling
                text = str(item) if item is not None else ""
                
                # Sanitize text based on font availability
                if not use_unicode_font:
                     replacements = {'−': '-', '–': '-', '—': '-'}
                     for old, new in replacements.items():
                         text = text.replace(old, new)
                     
                     # Final safety net: strictly encode to latin-1
                     text = text.encode('latin-1', 'replace').decode('latin-1')
                
                # Simple truncation to avoid breaking layout
                text = (text[:20] + '..') if len(text) > 20 else text
                pdf.cell(40, 10, text, border=1, new_x=XPos.RIGHT, new_y=YPos.TOP)
            pdf.ln()
            
        pdf.output(output_path)
    except Exception as e:
        raise RuntimeError(f"XLSX conversion failed: {e}")

def convert_image_to_pdf(input_path: str, output_path: str):
    try:
        with Image.open(input_path) as img:
            img.verify()
        
        with open(output_path, "wb") as f:
            f.write(img2pdf.convert(input_path))
    except Exception as e:
        raise RuntimeError(f"Image conversion failed: {e}")

def convert_pptx_to_pdf(input_path: str, output_path: str):
    # Try High Fidelity Conversion first
    if convert_with_libreoffice(input_path, output_path):
        return

    try:
        # Load PowerPoint presentation
        prs = Presentation(input_path)
        
        # Create PDF
        pdf = FPDF(orientation='L')  # Landscape for slides
        pdf.set_auto_page_break(auto=False)
        
        # Add Unicode font support (bundled Roboto)
        current_dir = os.path.dirname(os.path.abspath(__file__))
        font_path = os.path.join(current_dir, "fonts", "Roboto-Regular.ttf")
        font_name = "Roboto"
        use_unicode_font = False
        
        if os.path.exists(font_path):
            try:
                pdf.add_font(font_name, fname=font_path)
                use_unicode_font = True
            except Exception as e:
                print(f"Failed to load font {font_path}: {e}")
        else:
             print(f"Font not found at {font_path}")
        
        # Extract text from slides
        for slide in prs.slides:
            pdf.add_page()
            
            if use_unicode_font:
                pdf.set_font(font_name, size=12)
            else:
                pdf.set_font("Helvetica", size=12)
            
            y_position = 20
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        # Sanitize text based on font availability
                        if not use_unicode_font:
                            # Replace common problematic characters
                            replacements = {
                                '−': '-',  # minus sign
                                '–': '-',  # en dash
                                '—': '-',  # em dash
                                '"': '"',
                                '"': '"',
                                ''': "'",
                                ''': "'",
                                '…': '...'
                            }
                            for old, new in replacements.items():
                                text = text.replace(old, new)
                            
                            # Final safety net: strictly encode to latin-1, replacing errors with '?'
                            # This prevents the application from crashing
                            text = text.encode('latin-1', 'replace').decode('latin-1')

                        pdf.set_xy(20, y_position)
                        pdf.multi_cell(0, 10, text)
                        y_position += 15
        
        pdf.output(output_path)
    except Exception as e:
        raise RuntimeError(f"PowerPoint conversion failed: {e}")

def convert_html_to_pdf(input_path: str, output_path: str):
    # Try High Fidelity Conversion first (LibreOffice)
    if convert_with_libreoffice(input_path, output_path):
        return

    try:
        # Fallback: Use WeasyPrint (modern, reliable HTML/CSS support)
        from weasyprint import HTML, CSS
        
        # WeasyPrint needs a base_url so it can find images/css if relative
        # Since we just have a single HTML file here, we point it to the file's dir
        base = os.path.dirname(input_path)
        
        doc = HTML(filename=input_path, base_url=base)
        doc.write_pdf(output_path)
        
    except ImportError:
         print("WeasyPrint not found, falling back to basic conversion.")
         # Double Fallback: Simple text extraction if WeasyPrint is missing
         try:
            from reportlab.platypus import SimpleDocTemplate, Paragraph
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib.pagesizes import letter
            from xml.sax.saxutils import escape
            import re
            
            with open(input_path, 'r', encoding='utf-8') as f:
                html_content = f.read()

            pdf = SimpleDocTemplate(output_path, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            
            # Use regex to strip tags, but then ESCAPE content for Paragraph
            text = re.sub('<[^<]+?>', '', html_content)
            
            for line in text.split('\n'):
                if line.strip():
                    # Escape special characters so ReportLab doesn't crash on <, >, &
                    safe_text = escape(line.strip())
                    p = Paragraph(safe_text, styles['Normal'])
                    story.append(p)
            
            pdf.build(story)
         except Exception as e_inner:
             raise RuntimeError(f"Basic HTML conversion failed: {e_inner}")

    except Exception as e:
        raise RuntimeError(f"HTML conversion failed: {e}")

# PDF to Other Formats
def convert_pdf_to_jpg(input_path: str, output_path: str):
    try:
        # Convert PDF to images (first page)
        images = convert_from_path(input_path, first_page=1, last_page=1)
        
        if images:
            images[0].save(output_path, 'JPEG')
    except Exception as e:
        raise RuntimeError(f"PDF to JPG conversion failed: {e}")

def convert_pdf_to_docx(input_path: str, output_path: str):
    try:
        from pdf2docx import Converter
        
        # Use pdf2docx for structure-preserving conversion
        cv = Converter(input_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()
        
    except ImportError:
        # Fallback to text extraction if pdf2docx is missing
        try:
             pdf = PdfReader(input_path)
             doc = Document()
             for page in pdf.pages:
                 text = page.extract_text()
                 if text:
                     doc.add_paragraph(text)
                     doc.add_page_break()
             doc.save(output_path)
        except Exception as e_inner:
             raise RuntimeError(f"Basic PDF to Word conversion failed: {e_inner}")

    except Exception as e:
        raise RuntimeError(f"PDF to Word conversion failed: {e}")

def convert_pdf_to_xlsx(input_path: str, output_path: str):
    try:
        # Extract text from PDF
        with pdfplumber.open(input_path) as pdf:
            wb = openpyxl.Workbook()
            ws = wb.active
            
            for page in pdf.pages:
                # Try to extract tables
                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        for row in table:
                            ws.append(row)
                else:
                    # If no tables, extract text
                    text = page.extract_text()
                    if text:
                        ws.append([text])
            
            wb.save(output_path)
    except Exception as e:
        raise RuntimeError(f"PDF to Excel conversion failed: {e}")

def convert_pdf_to_pptx(input_path: str, output_path: str):
    try:
        from pptx.util import Inches
        
        # Extract text from PDF
        pdf = PdfReader(input_path)
        
        # Create PowerPoint
        prs = Presentation()
        
        for i, page in enumerate(pdf.pages):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            text = page.extract_text()
            if text:
                title.text = f"Page {i + 1}"
                content.text = text[:500]  # Limit text length
        
        prs.save(output_path)
    except Exception as e:
        raise RuntimeError(f"PDF to PowerPoint conversion failed: {e}")


def compress_pdf(input_path: str, output_path: str, compression_level: str = "medium"):
    """
    Compress a PDF file to reduce its size.
    
    Args:
        input_path: Path to the input PDF file
        output_path: Path to save the compressed PDF
        compression_level: 'low', 'medium', or 'high' compression
    """
    try:
        import fitz  # PyMuPDF
        
        # Open the PDF
        doc = fitz.open(input_path)
        
        # Set compression parameters based on level
        if compression_level == "low":
            # Light compression - preserve quality
            garbage = 1
            deflate = True
            image_quality = 95
        elif compression_level == "high":
            # Maximum compression - may reduce quality
            garbage = 4
            deflate = True
            image_quality = 60
        else:  # medium (default)
            garbage = 3
            deflate = True
            image_quality = 80
        
        # Process each page to compress images
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            # Get all images on the page
            image_list = page.get_images(full=True)
            
            for img_index, img in enumerate(image_list):
                xref = img[0]
                
                try:
                    # Extract the image
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    # Only compress JPEG and PNG images
                    if image_ext.lower() in ['jpeg', 'jpg', 'png']:
                        from PIL import Image
                        import io
                        
                        # Open and compress the image
                        pil_image = Image.open(io.BytesIO(image_bytes))
                        
                        # Convert to RGB if necessary (for JPEG)
                        if pil_image.mode in ('RGBA', 'P'):
                            pil_image = pil_image.convert('RGB')
                        
                        # Reduce resolution for high compression
                        if compression_level == "high":
                            max_size = 1024
                            if pil_image.width > max_size or pil_image.height > max_size:
                                ratio = min(max_size / pil_image.width, max_size / pil_image.height)
                                new_size = (int(pil_image.width * ratio), int(pil_image.height * ratio))
                                pil_image = pil_image.resize(new_size, Image.LANCZOS)
                        
                        # Compress to JPEG
                        img_buffer = io.BytesIO()
                        pil_image.save(img_buffer, format='JPEG', quality=image_quality, optimize=True)
                        compressed_bytes = img_buffer.getvalue()
                        
                        # Only replace if the compressed image is smaller
                        if len(compressed_bytes) < len(image_bytes):
                            # Update the image in the PDF
                            page.replace_image(xref, filename=None, stream=compressed_bytes)
                except Exception as img_error:
                    # Skip problematic images and continue
                    print(f"Warning: Could not compress image {img_index} on page {page_num}: {img_error}")
                    continue
        
        # Save with compression options
        doc.save(
            output_path,
            garbage=garbage,
            deflate=deflate,
            deflate_images=True,
            deflate_fonts=True,
            clean=True
        )
        
        doc.close()
        
    except ImportError:
        # Fallback using pypdf if fitz is not available
        try:
            from pypdf import PdfReader, PdfWriter
            
            reader = PdfReader(input_path)
            writer = PdfWriter()
            
            for page in reader.pages:
                page.compress_content_streams()
                writer.add_page(page)
            
            # Remove duplicate objects
            writer.add_metadata(reader.metadata or {})
            
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)
                
        except Exception as e_inner:
            raise RuntimeError(f"PDF compression failed (fallback): {e_inner}")
    
    except Exception as e:
        raise RuntimeError(f"PDF compression failed: {e}")

